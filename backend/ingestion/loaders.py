import os
from bs4 import BeautifulSoup
from docx import Document


def load_document(file_path: str, mime_type: str) -> str:
    """
    Loads text from a file based on its mime_type/extension.
    Returns the extracted text.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = file_path.split('.')[-1].lower()

    if ext == 'txt' or mime_type == 'text/plain':
        return _load_txt(file_path)
    elif ext == 'pdf' or mime_type == 'application/pdf':
        return _load_pdf(file_path)
    elif ext in ['htm', 'html'] or mime_type == 'text/html':
        return _load_html(file_path)
    elif ext == 'docx' or mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return _load_docx(file_path)
    elif ext == 'pptx':
        return _load_pptx(file_path)
    elif ext in ['csv', 'xlsx']:
        return _load_spreadsheet(file_path, ext)
    elif ext == 'md':
        return _load_txt(file_path)  # Markdown is plaintext
    else:
        raise ValueError(f"Unsupported file format: {ext}")

def _load_txt(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def _load_pdf(file_path: str) -> str:
    import fitz # PyMuPDF
    doc = fitz.open(file_path)
    text = []
    for page in doc:
        page_text = page.get_text()
        text.append(f"\n[PAGE {page.number + 1}]\n{page_text}")
    doc.close()
    return "\n".join(text)

def _load_html(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8') as f:
        html_content = f.read()
    soup = BeautifulSoup(html_content, 'html.parser')
    return soup.get_text(separator='\n')

def _load_docx(file_path: str) -> str:
    doc = Document(file_path)
    return '\n'.join([para.text for para in doc.paragraphs])

def _load_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text = []
    for slide_num, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        text.append(f"\n[SLIDE {slide_num + 1}]\n" + "\n".join(slide_text))
    return "\n".join(text)

def _load_spreadsheet(file_path: str, ext: str) -> str:
    import pandas as pd
    if ext == 'csv':
        df = pd.read_csv(file_path)
    else:
        df = pd.read_excel(file_path)
    return df.to_markdown(index=False)
